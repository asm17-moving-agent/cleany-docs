#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Iterable, Sequence

sys.dont_write_bytecode = True

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf"}
EXCLUDED_DIRS = {".git", ".venv", "node_modules", "__pycache__", ".obsidian"}


def yaml_quote(value: object) -> str:
    return json.dumps("" if value is None else str(value), ensure_ascii=False)


def normalize_text(value: object) -> str:
    if value is None:
        return ""
    text = str(value)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    return text.strip()


def escape_cell(value: object) -> str:
    text = normalize_text(value)
    text = text.replace("|", "\\|")
    text = text.replace("\n", "<br>")
    return text


def trim_table(rows: Sequence[Sequence[object]]) -> list[list[str]]:
    normalized = [[escape_cell(cell) for cell in row] for row in rows]
    while normalized and all(cell == "" for cell in normalized[-1]):
        normalized.pop()
    if not normalized:
        return []
    max_cols = max((len(row) for row in normalized), default=0)
    padded = [row + [""] * (max_cols - len(row)) for row in normalized]
    while padded and max_cols > 0 and all(row[max_cols - 1] == "" for row in padded):
        max_cols -= 1
    return [row[:max_cols] for row in padded if max_cols > 0]


def table_to_markdown(rows: Sequence[Sequence[object]]) -> str:
    table = trim_table(rows)
    if not table:
        return ""
    header = table[0]
    body = table[1:] or [[""] * len(header)]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * len(header)) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def frontmatter(path: Path, source_format: str) -> str:
    return "\n".join(
        [
            "---",
            "type: raw-converted",
            f"source_file: {yaml_quote(path.name)}",
            f"source_format: {yaml_quote(source_format)}",
            "converter: office-to-markdown",
            "reviewers:",
            "  -",
            "---",
            "",
        ]
    )


def heading_for_docx_paragraph(style_name: str) -> int | None:
    match = re.search(r"Heading\s+(\d+)", style_name, flags=re.IGNORECASE)
    if match:
        return min(int(match.group(1)) + 1, 6)
    if style_name.lower() == "title":
        return 2
    return None


def docx_xml_text(element) -> str:
    parts: list[str] = []
    for node in element.iter():
        local = node.tag.rsplit("}", 1)[-1]
        if local == "t" and node.text:
            parts.append(node.text)
        elif local == "tab":
            parts.append("\t")
        elif local == "br":
            parts.append("\n")
    return normalize_text("".join(parts))


def docx_cell_text(cell) -> str:
    paragraphs: list[str] = []
    for node in cell._tc.iter():
        if node.tag.rsplit("}", 1)[-1] == "p":
            text = docx_xml_text(node)
            if text:
                paragraphs.append(text)
    return "\n".join(paragraphs)


def docx_table_rows(table) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in table.rows:
        seen_tc: set[int] = set()
        values: list[str] = []
        for cell in row.cells:
            tc_id = id(cell._tc)
            if tc_id in seen_tc:
                values.append("")
            else:
                seen_tc.add(tc_id)
                values.append(docx_cell_text(cell))
        rows.append(values)
    return rows


def iter_docx_blocks(document):
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def convert_docx(path: Path) -> str:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(path)
    lines: list[str] = [f"# {path.stem}", ""]

    for block in iter_docx_blocks(doc):
        if isinstance(block, Paragraph):
            text = docx_xml_text(block._p)
            if not text:
                continue
            style_name = block.style.name if block.style is not None else ""
            heading_level = heading_for_docx_paragraph(style_name)
            lower_style = style_name.lower()
            if heading_level:
                lines.extend(["#" * heading_level + " " + text, ""])
            elif "list bullet" in lower_style:
                lines.extend([f"- {text}", ""])
            elif "list number" in lower_style:
                lines.extend([f"1. {text}", ""])
            else:
                lines.extend([text, ""])
        elif isinstance(block, Table):
            md = table_to_markdown(docx_table_rows(block))
            if md:
                lines.extend([md, ""])

    return "\n".join(lines).rstrip() + "\n"


def convert_xlsx(path: Path, max_rows: int | None = None) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = [f"# {path.stem}", ""]
    for sheet in wb.worksheets:
        lines.extend([f"## 시트: {sheet.title}", ""])
        rows: list[list[object]] = []
        for idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if max_rows is not None and idx > max_rows:
                break
            rows.append(list(row))
        md = table_to_markdown(rows)
        if md:
            lines.extend([md, ""])
        else:
            lines.extend(["_빈 시트_", ""])
        if max_rows is not None and sheet.max_row and sheet.max_row > max_rows:
            lines.extend([f"> {max_rows}행까지만 변환됨. 전체 행 수: {sheet.max_row}", ""])
    wb.close()
    return "\n".join(lines).rstrip() + "\n"


def shape_text_lines(shape) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    result: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = normalize_text("".join(run.text for run in paragraph.runs) or paragraph.text)
        if not text:
            continue
        level = getattr(paragraph, "level", 0) or 0
        if level > 0:
            result.append("  " * level + f"- {text}")
        else:
            result.append(text)
    return result


def convert_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    lines: list[str] = [f"# {path.stem}", ""]
    for slide_no, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title = normalize_text(title_shape.text) if title_shape is not None else ""
        if title:
            lines.extend([f"## 슬라이드 {slide_no}: {title}", ""])
        else:
            lines.extend([f"## 슬라이드 {slide_no}", ""])

        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                md = table_to_markdown(rows)
                if md:
                    lines.extend([md, ""])
            else:
                text_lines = shape_text_lines(shape)
                if text_lines:
                    lines.extend(text_lines + [""])
    return "\n".join(lines).rstrip() + "\n"


def convert_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    lines: list[str] = [f"# {path.stem}", ""]
    for page_no, page in enumerate(reader.pages, start=1):
        lines.extend([f"## 페이지 {page_no}", ""])
        text = normalize_text(page.extract_text() or "")
        if text:
            lines.extend([text, ""])
        else:
            lines.extend(["_텍스트를 추출하지 못함_", ""])
    return "\n".join(lines).rstrip() + "\n"


def convert_file(path: Path, include_frontmatter: bool, max_rows: int | None) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        body = convert_docx(path)
    elif suffix == ".xlsx":
        body = convert_xlsx(path, max_rows=max_rows)
    elif suffix == ".pptx":
        body = convert_pptx(path)
    elif suffix == ".pdf":
        body = convert_pdf(path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {suffix}")
    if include_frontmatter:
        return frontmatter(path, suffix.lstrip(".")) + body
    return body


def iter_supported_files(root: Path) -> Iterable[Path]:
    for path in root.rglob("*"):
        if any(part in EXCLUDED_DIRS for part in path.parts):
            continue
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def output_path_for(input_root: Path, output_root: Path, file_path: Path) -> Path:
    rel = file_path.relative_to(input_root)
    return output_root / rel.with_suffix(".md")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Office/PDF 파일을 deterministic Markdown으로 변환한다.")
    parser.add_argument("input", help="입력 파일 또는 폴더")
    parser.add_argument("--output", "-o", help="출력 파일 또는 폴더. 생략 시 단일 파일은 stdout 출력")
    parser.add_argument("--no-frontmatter", action="store_true", help="YAML frontmatter를 출력하지 않음")
    parser.add_argument("--max-rows", type=int, default=None, help="XLSX 시트별 최대 변환 행 수")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_path = Path(args.input).resolve()
    include_frontmatter = not args.no_frontmatter

    if not input_path.exists():
        print(f"입력 경로가 존재하지 않음: {input_path}", file=sys.stderr)
        return 1

    if input_path.is_dir():
        if not args.output:
            print("폴더 변환에는 --output 출력 폴더가 필요함", file=sys.stderr)
            return 1
        output_root = Path(args.output).resolve()
        converted = 0
        for file_path in sorted(iter_supported_files(input_path)):
            out = output_path_for(input_path, output_root, file_path)
            out.parent.mkdir(parents=True, exist_ok=True)
            out.write_text(convert_file(file_path, include_frontmatter, args.max_rows), encoding="utf-8")
            print(f"converted: {file_path} -> {out}")
            converted += 1
        print(f"총 {converted}개 파일 변환 완료")
        return 0

    if input_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        print(f"지원하지 않는 파일 형식: {input_path.suffix}", file=sys.stderr)
        return 1

    markdown = convert_file(input_path, include_frontmatter, args.max_rows)
    if args.output:
        output_path = Path(args.output).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(markdown, encoding="utf-8")
        print(f"converted: {input_path} -> {output_path}")
    else:
        print(markdown, end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
